from flask import Flask, request, redirect, url_for, render_template, send_from_directory, flash
import os
from werkzeug.utils import secure_filename
from gtts import gTTS
from pydub import AudioSegment
import fitz
from langchain.schema import Document
import chromadb
import chromadb.utils.embedding_functions as embedding_functions
import google.generativeai as genai
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
from pptx import Presentation
from bs4 import BeautifulSoup

apikey=''
genai.configure(api_key=apikey)
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['ALLOWED_EXTENSIONS'] = {'txt', 'pdf','pptx','html'}
app.config['MAX_CONTENT_LENGTH'] = 20 * 1024 * 1024  # 20 MB limit

checkpoint = "sshleifer/distilbart-cnn-12-6"
tokenizer_s = AutoTokenizer.from_pretrained(checkpoint)
model_s = AutoModelForSeq2SeqLM.from_pretrained(checkpoint)

if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

@app.route('/', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        if 'files[]' not in request.files:
            flash('No file part')
            return redirect(request.url)
        files = request.files.getlist('files[]')
        for file in files:
            if file and allowed_file(file.filename):
                filename = secure_filename(file.filename)
                file.save(os.path.join(app.config['UPLOAD_FOLDER'], filename))
        
        # Process files after uploading
        extracted_data = load_pdf(app.config['UPLOAD_FOLDER'])
        text_chunks = []
        for doc in extracted_data:
            chunks = split_text_recursively(doc.page_content, max_length=500, chunk_overlap=20)
            for chunk in chunks:
                text_chunks.append(Document(page_content=chunk, metadata=doc.metadata))
        client = chromadb.PersistentClient(path='embeddings/gemini')
        client.delete_collection('pdf_rag')
        google_ef = embedding_functions.GoogleGenerativeAiEmbeddingFunction(api_key=apikey)
        collection = client.get_or_create_collection(name='pdf_rag', embedding_function=google_ef)
        collection.add(documents=[d.page_content for d in text_chunks], 
                       metadatas=[d.metadata for d in text_chunks], 
                       ids=[str(i) for i in range(len(text_chunks))])
        
        return redirect(url_for('uploaded_files'))
    
    return render_template('index.html')
@app.route('/uploads/<filename>')
def uploaded_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

@app.route('/uploaded_files', methods=['GET', 'POST'])
def uploaded_files():
    files = os.listdir(app.config['UPLOAD_FOLDER'])
    if request.method == 'POST' :
        if 'process' in request.form:
            # Call the function to process files
            extracted_data = load_pdf(app.config['UPLOAD_FOLDER'])
            text_chunks = []
            for doc in extracted_data:
                chunks = split_text_recursively(doc.page_content, max_length=500, chunk_overlap=20)
                for chunk in chunks:
                    text_chunks.append(Document(page_content=chunk, metadata=doc.metadata))
            client = chromadb.PersistentClient(path='embeddings/gemini')
            client.delete_collection('pdf_rag')
            google_ef=embedding_functions.GoogleGenerativeAiEmbeddingFunction(api_key=apikey)
            collection = client.get_or_create_collection(name='pdf_rag', embedding_function=google_ef)
            collection.add(documents=[d.page_content for d in text_chunks], 
                           metadatas=[d.metadata for d in text_chunks], 
                           ids=[str(i) for i in range(len(text_chunks))])
        elif 'ask' in request.form:
            query = request.form['query']
            return redirect(url_for('answer', query=query))
        elif 'summary' in request.form:
            selected_file = request.form.get('file')
            return redirect(url_for('summary', filename=selected_file))
    return render_template('uploaded_files.html', files=files)

@app.route('/delete/<filename>', methods=['POST'])
def delete_file(filename):
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    if os.path.exists(file_path):
        os.remove(file_path)

    # Process files after deleting
    extracted_data = load_pdf(app.config['UPLOAD_FOLDER'])
    text_chunks = []
    for doc in extracted_data:
        chunks = split_text_recursively(doc.page_content, max_length=500, chunk_overlap=20)
        for chunk in chunks:
            text_chunks.append(Document(page_content=chunk, metadata=doc.metadata))
    client = chromadb.PersistentClient(path='embeddings/gemini')
    client.delete_collection('pdf_rag')
    google_ef = embedding_functions.GoogleGenerativeAiEmbeddingFunction(api_key=apikey)
    collection = client.get_or_create_collection(name='pdf_rag', embedding_function=google_ef)
    collection.add(documents=[d.page_content for d in text_chunks],
                   metadatas=[d.metadata for d in text_chunks],
                   ids=[str(i) for i in range(len(text_chunks))])

    return redirect(url_for('uploaded_files'))


@app.route('/summary/<filename>', methods=['GET', 'POST'])
def summary(filename):
    if request.method == 'POST':
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        if filename.endswith('.txt'):
            with open(file_path, 'r') as f:
                text = f.read()
        elif filename.endswith('.pdf'):
            text = extract_text_from_pdf(file_path)
        elif filename.endswith('.pptx'):
            presentation = Presentation(file_path)
            text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            text='\n'.join(text)
        elif filename.endswith('.html'):
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            soup = BeautifulSoup(html_content, 'html.parser')
            text = soup.get_text()
        else:
            flash('Unsupported file type')
            return redirect(url_for('upload_file'))
        
        summary_text = summarize_text(text)
        tts = gTTS(summary_text, lang='en')
        x=filename.split('.')[0]
        print(x)
        x=filename+"_summary.mp3"
        summary_audio_path = os.path.join(app.config['UPLOAD_FOLDER'], x)
        tts.save(summary_audio_path)
        return render_template('summary.html', filename=filename, summary=summary_text, audio_path=summary_audio_path)
    return redirect(url_for('uploaded_files'))

@app.route('/answer')
def answer():
    query = request.args.get('query')
    client = chromadb.PersistentClient(path='embeddings/gemini')
    google_ef=embedding_functions.GoogleGenerativeAiEmbeddingFunction(api_key=apikey)
    collection = client.get_or_create_collection(name='pdf_rag', embedding_function=google_ef)
    results, sources = find_relevant_context(query, collection, 3)
    prompt, sources = create_prompt_for_gemini(query, results, sources)
    answer_text = generate_answer_from_gemini(prompt)
    return render_template('answer.html', query=query, answer=answer_text.text, sources=sources)

def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def summarize_text(text):
    sentences = text.split('.')
    length = 0
    chunk = ""
    chunks = []
    count = -1
    for sentence in sentences:
        count += 1
        combined_length = len(tokenizer_s.tokenize(sentence)) + length
        if combined_length <= tokenizer_s.model_max_length:
            chunk += sentence + ". "
            length = combined_length
            if count == len(sentences) - 1:
                chunks.append(chunk.strip())
        else:
            chunks.append(chunk.strip())
            length = len(tokenizer_s.tokenize(sentence))
            chunk = sentence + ". "
    summaries = []
    for chunk in chunks:
        inputs = tokenizer_s(chunk, return_tensors="pt", max_length=tokenizer_s.model_max_length, truncation=True)
        summary_ids = model_s.generate(inputs["input_ids"], max_length=150, min_length=30, length_penalty=2.0, num_beams=4, early_stopping=True)
        summary = tokenizer_s.decode(summary_ids[0], skip_special_tokens=True)
        summaries.append(summary)
    return " ".join(summaries)

def load_pdf(directory):
    documents = []
    for filename in os.listdir(directory):
        file_path = os.path.join(directory, filename)
        try:
            if filename.endswith(".pdf"):
                doc = fitz.open(file_path)
                text = ""
                for page in doc:
                    text += page.get_text()
                documents.append(Document(page_content=text, metadata={"source": file_path}))
            elif filename.endswith(".txt"):
                with open(file_path, 'r', encoding='utf-8') as file:
                    text = file.read()
                    documents.append(Document(page_content=text, metadata={"source": file_path}))
            elif filename.endswith('.pptx'):
                presentation = Presentation(file_path)
                text = []
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                text='\n'.join(text)
                documents.append(Document(page_content=text, metadata={"source": file_path}))
            elif filename.endswith('.html'):
                with open(file_path, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                soup = BeautifulSoup(html_content, 'html.parser')
                text = soup.get_text()
                documents.append(Document(page_content=text, metadata={"source": file_path}))
        except Exception as e:
            print("Error loading :",filename,e)
    return documents

def split_text_recursively(text, max_length=500, chunk_overlap=20):
    chunks = []
    start = 0
    while start < len(text):
        end = start + max_length
        if end < len(text):
            end = text.rfind(' ', start, end) + 1
            if end <= start:
                end = start + max_length
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end - chunk_overlap
        if start >= len(text):
            break
    return chunks

def build_escaped_context(context):
    escaped_context = ''
    sources = []
    for item in context:
        escaped_context += item['document'] + "\n"
        sources.append(item['metadata'])
    return escaped_context, sources

def find_relevant_context(query, db, n_result=3):
    results = db.query(query_texts=[query], n_results=n_result)
    documents = results['documents'][0]
    metadatas = results['metadatas'][0]
    context = [{'document': doc, 'metadata': meta} for doc, meta in zip(documents, metadatas)]
    escaped_context, source = build_escaped_context(context)
    return escaped_context, source

def create_prompt_for_gemini(query, context, sources):
    prompt = f"""
    You are an agent that answers questions using the text from the context below.
    Both the question and context is shared with you and you should answer the question on the basis of the context and not hallucinate.
    If the context does not have enough information for you, inform about the absence of relevant context as part of your answer.

    Context : {context}

    Question : {query}

    Answer :
    """
    return prompt, sources


def generate_answer_from_gemini(prompt):
    model = genai.GenerativeModel('gemini-1.5-pro-latest')
    result = model.generate_content(prompt)
    return result

if __name__ == '__main__':
    app.secret_key = 'supersecretkey'
    app.run(debug=True, use_reloader=False)
